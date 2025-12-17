import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import io
from pptx import Presentation
from pptx.util import Inches

# =============================================================================
# PAGE CONFIG
# =============================================================================
st.set_page_config(
    page_title="Oil & Gas Analytics Dashboard",
    page_icon="🛢️",
    layout="wide"
)

# =============================================================================
# HELPERS
# =============================================================================

def find_column(df, level0, level1_contains=""):
    for col in df.columns:
        if (
            str(col[0]).strip().upper() == level0.upper()
            and level1_contains.upper() in str(col[1]).replace("\n", " ").upper()
        ):
            return col
    return None


def normalize_numeric(series):
    return (
        series.astype(str)
        .str.replace(",", "", regex=False)
        .replace("", np.nan)
        .pipe(pd.to_numeric, errors="coerce")
    )

# =============================================================================
# DATA EXTRACTION
# =============================================================================

@st.cache_data(show_spinner=False)
def extract_data(file):
    df = pd.read_excel(
        file,
        sheet_name="Report",
        skiprows=6,
        header=[0, 1]
    )

    field = find_column(df, "Field")
    well = find_column(df, "RUNNING WELLS")
    net_bo = find_column(df, "TOTAL PRODUCTION", "Net BO")
    net_diff = find_column(df, "TOTAL PRODUCTION", "Net diff")
    wc = find_column(df, "W/C", "%")

    if None in [field, well, net_bo, net_diff]:
        return None, None

    df[net_bo] = normalize_numeric(df[net_bo])
    df[net_diff] = normalize_numeric(df[net_diff])
    if wc:
        df[wc] = normalize_numeric(df[wc])

    # Stop at TOTAL row
    total_mask = df[field].astype(str).str.contains("TOTAL", na=False)
    if total_mask.any():
        df = df.loc[: total_mask.idxmax() - 1]

    # Clean wells
    df = df[
        df[well].notna() &
        (df[well].astype(str).str.strip() != "")
    ]

    cols = {
        "field": field,
        "well": well,
        "net_bo": net_bo,
        "net_diff": net_diff,
        "wc": wc
    }

    return df.reset_index(drop=True), cols

# =============================================================================
# AI ANOMALY DETECTION (FIXED)
# =============================================================================

def detect_anomalies(df, cols):
    df = df.copy()

    # --- Z-score (Net Diff BO) ---
    std = df[cols["net_diff"]].std()
    if std == 0 or pd.isna(std):
        df["z_net_diff"] = 0
    else:
        df["z_net_diff"] = (
            (df[cols["net_diff"]] - df[cols["net_diff"]].mean()) / std
        )

    # --- IQR (Net BO) ---
    q1 = df[cols["net_bo"]].quantile(0.25)
    q3 = df[cols["net_bo"]].quantile(0.75)
    iqr = q3 - q1

    lower_bo = q1 - 1.5 * iqr
    upper_bo = q3 + 1.5 * iqr

    # --- Flags ---
    df["FLAG_EXTREME_DIFF"] = df["z_net_diff"].abs() > 3
    df["FLAG_ZERO_BO"] = df[cols["net_bo"]] == 0
    df["FLAG_IQR_BO"] = (df[cols["net_bo"]] < lower_bo) | (df[cols["net_bo"]] > upper_bo)
    df["FLAG_HIGH_WC"] = df[cols["wc"]] > 75 if cols["wc"] else False

    df["ANOMALY"] = (
        df["FLAG_EXTREME_DIFF"] |
        df["FLAG_ZERO_BO"] |
        df["FLAG_IQR_BO"] |
        df["FLAG_HIGH_WC"]
    )

    # --- Human-readable reasons (SAFE) ---
    def build_reason(row):
        reasons = []
        if row["FLAG_EXTREME_DIFF"]:
            reasons.append("Extreme Net Diff BO (Z-score > 3)")
        if row["FLAG_ZERO_BO"]:
            reasons.append("Zero Net BO")
        if row["FLAG_IQR_BO"]:
            reasons.append("Abnormal Net BO (IQR outlier)")
        if row["FLAG_HIGH_WC"]:
            reasons.append("High Water Cut (>75%)")
        return "; ".join(reasons)

    df["REASON"] = df.apply(build_reason, axis=1)

    return df[df["ANOMALY"]].reset_index(drop=True)

# =============================================================================
# ADVANCED VISUALS
# =============================================================================

def create_visuals(df, cols):
    fig, axes = plt.subplots(1, 3, figsize=(32, 10))
    fig.set_dpi(300)

    # 1️⃣ Top Net Diff BO
    diff_df = (
        df.assign(abs_diff=df[cols["net_diff"]].abs())
        .sort_values("abs_diff", ascending=False)
        .head(15)
    )

    colors = [
        "orange" if bo == 0 else ("green" if v > 0 else "red")
        for v, bo in zip(diff_df[cols["net_diff"]], diff_df[cols["net_bo"]])
    ]

    axes[0].bar(diff_df[cols["well"]], diff_df[cols["net_diff"]], color=colors)
    axes[0].axhline(0, color="black", linewidth=2)
    axes[0].set_title("Top 15 Net Diff BO Wells")
    axes[0].tick_params(axis="x", rotation=45)

    # 2️⃣ Top W/C (exclude zero BO)
    if cols["wc"]:
        wc_df = df[df[cols["net_bo"]] > 0].nlargest(10, cols["wc"])
        axes[1].barh(wc_df[cols["well"]], wc_df[cols["wc"]])
        axes[1].set_title("Top 10 Wells by W/C")

    # 3️⃣ Top Net BO
    bo_df = df.nlargest(10, cols["net_bo"])
    axes[2].barh(bo_df[cols["well"]], bo_df[cols["net_bo"]])
    axes[2].set_title("Top 10 Net BO Wells")

    plt.tight_layout()
    return fig

# =============================================================================
# POWERPOINT EXPORT
# =============================================================================

def create_ppt(df, anomalies, fig, cols):
    prs = Presentation()

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Production Analysis Report"
    slide.placeholders[1].text = "Automated AI-Powered Dashboard"

    # AI Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Anomaly Summary"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    if anomalies.empty:
        tf.text = "No anomalies detected."
    else:
        for _, r in anomalies.iterrows():
            tf.add_paragraph().text = f"{r[cols['well']]} → {r['REASON']}"

    # Visuals
    img = io.BytesIO()
    fig.savefig(img, format="png", dpi=300)
    img.seek(0)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(img, Inches(0.5), Inches(0.5), width=Inches(9))

    ppt = io.BytesIO()
    prs.save(ppt)
    ppt.seek(0)
    return ppt

# =============================================================================
# UI
# =============================================================================

st.title("🛢️ Oil & Gas Analytics Dashboard")

uploaded_file = st.file_uploader(
    "Upload Production Excel File",
    type=["xlsx", "xlsm"]
)

if uploaded_file:
    with st.spinner("Processing Excel file..."):
        df, cols = extract_data(uploaded_file)

    if df is None:
        st.error("❌ Required columns not found in Excel file.")
        st.stop()

    anomalies = detect_anomalies(df, cols)
    fig = create_visuals(df, cols)

    # KPIs
    st.subheader("📊 Key Performance Indicators")
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Wells", len(df))
    c2.metric("Anomalies", len(anomalies))
    c3.metric("Zero Net BO", int((df[cols["net_bo"]] == 0).sum()))
    c4.metric("High W/C", int((df[cols["wc"]] > 75).sum()) if cols["wc"] else 0)

    # Anomalies table
    st.subheader("🚨 AI-Detected Anomalies")
    st.dataframe(
        anomalies[[cols["well"], cols["net_bo"], cols["net_diff"], cols["wc"], "REASON"]]
        if not anomalies.empty else anomalies,
        use_container_width=True
    )

    # Visuals
    st.subheader("📈 Advanced Analytics")
    st.pyplot(fig)
    plt.close(fig)

    # PPT
    ppt = create_ppt(df, anomalies, fig, cols)
    st.download_button(
        "📥 Download PowerPoint Report",
        ppt,
        "production_ai_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

else:
    st.info("👆 Upload your production Excel file to start analysis")
